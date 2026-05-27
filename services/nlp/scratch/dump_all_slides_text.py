import os
import sys
from pptx import Presentation

sys.stdout.reconfigure(encoding='utf-8')
path = r"C:\Projects\Samtool\OneDrive_1_5-22-2026\12-Automotive-Manufacturing.pptx"
prs = Presentation(path)
for idx, slide in enumerate(prs.slides):
    print(f"--- Slide {idx+1} ---")
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                print(text)
