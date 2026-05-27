import os
from pptx import Presentation

FOLDER = r"C:\Projects\Samtool\OneDrive_1_5-22-2026"
files = [f for f in os.listdir(FOLDER) if f.endswith(".pptx")]

for file in sorted(files):
    path = os.path.join(FOLDER, file)
    try:
        prs = Presentation(path)
        print(f"=== {file} ===")
        # Print text of slide 1 and 2
        for idx in range(min(len(prs.slides), 3)):
            print(f"  Slide {idx+1}:")
            for shape in prs.slides[idx].shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        # Print first 3 lines
                        lines = [line.strip() for line in text.split('\n') if line.strip()]
                        for line in lines[:3]:
                            print(f"    {line}")
    except Exception as e:
        print(f"Error reading {file}: {e}")
