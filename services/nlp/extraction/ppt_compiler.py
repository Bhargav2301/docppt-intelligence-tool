import io
import logging
from pptx import Presentation
from typing import Dict, Any, List

logger = logging.getLogger(__name__)

def compile_ppt(original_file_bytes: bytes, modifications: List[Dict[str, Any]]) -> bytes:
    """
    Takes an original PPTX file and a list of modifications, overwriting text at the precise run level.
    Modifications format: 
    [{"slide_index": 0, "shape_id": "title", "paragraph_index": 0, "run_index": 0, "new_text": "Cleaned Title"}]
    Returns the compiled PPTX as binary bytes.
    """
    try:
        prs = Presentation(io.BytesIO(original_file_bytes))
    except Exception as e:
        raise ValueError(f"Failed to load original PPTX for compilation: {str(e)}")

    success_count = 0
    skip_count = 0

    for mod in modifications:
        slide_idx = mod.get("slide_index")
        shape_id = str(mod.get("shape_id"))
        para_idx = mod.get("paragraph_index")
        run_idx = mod.get("run_index")
        new_text = mod.get("new_text", "")

        # Graceful failure block to handle index drift if the file was mutated
        try:
            slide = prs.slides[slide_idx]
            
            # Find shape by id
            target_shape = None
            for shape in slide.shapes:
                if str(shape.shape_id) == shape_id:
                    target_shape = shape
                    break
                    
            if not target_shape:
                raise KeyError(f"Shape ID {shape_id} not found on slide {slide_idx}")
                
            if not hasattr(target_shape, "text_frame") or not target_shape.text_frame:
                raise ValueError(f"Shape {shape_id} has no text frame")
                
            paragraph = target_shape.text_frame.paragraphs[para_idx]
            run = paragraph.runs[run_idx]
            
            # Perform surgical overwrite
            run.text = new_text
            success_count += 1
            
        except (IndexError, KeyError, ValueError, AttributeError) as e:
            logger.warning(
                f"Skipping modification at Slide {slide_idx}, Shape {shape_id}, "
                f"Para {para_idx}, Run {run_idx}. Reason: {str(e)}"
            )
            skip_count += 1
            continue

    logger.info(f"Compilation finished. Successfully applied {success_count} edits. Skipped {skip_count} due to index drift.")

    # Save to in-memory buffer
    output_buffer = io.BytesIO()
    prs.save(output_buffer)
    return output_buffer.getvalue()
