import io
from pptx import Presentation
from typing import Dict, Any, List

def extract_ppt_text(file_source) -> Dict[str, Any]:
    """
    Parses a PPTX file from bytes or file path and extracts all text down to the run level.
    Returns a dictionary mapping exactly to the ppt_text_segments schema requirements.
    """
    try:
        if isinstance(file_source, bytes):
            prs = Presentation(io.BytesIO(file_source))
        else:
            prs = Presentation(file_source)
    except Exception as e:
        raise ValueError(f"Failed to parse PPTX file: {str(e)}")

    slides_data = []
    total_segments = 0

    for slide_idx, slide in enumerate(prs.slides):
        segments = []
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame") or not shape.text_frame:
                continue
                
            for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                for run_idx, run in enumerate(paragraph.runs):
                    text = run.text.strip()
                    if text:
                        shape_name = str(getattr(shape, "name", "")).lower()
                        if "title" in shape_name or "header" in shape_name:
                            role = "title"
                        elif "note" in shape_name or "speaker" in shape_name:
                            role = "speaker_note"
                        elif getattr(paragraph, "level", 0) > 0 or len(text.split()) < 15:
                            role = "bullet"
                        else:
                            role = "body"
                            
                        segments.append({
                            "shape_id": str(shape.shape_id),
                            "paragraph_index": para_idx,
                            "run_index": run_idx,
                            "original_text": text,
                            "normalized_text": text,
                            "role": role,
                            "flags": []
                        })
                        total_segments += 1
                        
        slides_data.append({
            "slide_index": slide_idx,
            "slide_number": slide_idx + 1,
            "segments": segments
        })

    if total_segments == 0:
        raise ValueError("No extractable text found. The presentation might be empty or image-only.")

    return {
        "total_slides": len(prs.slides),
        "total_segments": total_segments,
        "slides": slides_data
    }
