from pptx import Presentation

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "AI Strategy 2026 [Artifact: 123]"
subtitle.text = "Prepared by ChatGPT. Contact: support@openai.com [Draft v1.0]"

# Add a bullet slide
bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
title_shape.text = "Key Objectives"

body_shape = shapes.placeholders[1]
tf = body_shape.text_frame
tf.text = "Integrate AI pipelines."

p = tf.add_paragraph()
p.text = "Eliminate manual errors (Ref: [Doc-45])."
p.level = 1

p = tf.add_paragraph()
p.text = "Optimize for speed and cost."
p.level = 1

prs.save("sample_presentation.pptx")
print("Created sample_presentation.pptx")
