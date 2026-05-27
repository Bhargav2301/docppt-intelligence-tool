import os
import re
from pptx import Presentation

FOLDER = r"C:\Projects\Samtool\OneDrive_1_5-22-2026"
files = [f for f in os.listdir(FOLDER) if f.endswith(".pptx")]

ai_indicators = [
    r'\bkey\s+takeaways\b',
    r'\bin\s+conclusion\b',
    r'\boverview\b',
    r'\bseamless\b',
    r'\bscalable\b',
    r'\bleverage\b',
    r'\brobust\b',
    r'\bstreamline\b',
    r'\boptimiz\w+\b',
    r'\bsynerg\w+\b',
    r'\btransform\b',
    r'\bvalue\s+proposition\b',
    r'\bstrategic\b'
]

results = []

for file in sorted(files):
    path = os.path.join(FOLDER, file)
    try:
        prs = Presentation(path)
        slide_count = len(prs.slides)
        
        # Extract text and analyze
        all_text = []
        placeholder_count = 0
        ai_pattern_matches = 0
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            all_text.append(text)
                            
        full_text = " ".join(all_text)
        words = full_text.split()
        word_count = len(words)
        avg_words = round(word_count / slide_count, 1) if slide_count > 0 else 0
        
        # Check for placeholders
        placeholders = re.findall(r'\[\s*insert\s*.*?\]|\{\s*insert\s*.*?\}|<\s*insert\s*.*?>|\bTODO\b|\bTBD\b|\bplaceholder\b', full_text, re.IGNORECASE)
        placeholder_count = len(placeholders)
        
        # Check for AI indicator terms
        for pat in ai_indicators:
            matches = re.findall(pat, full_text, re.IGNORECASE)
            ai_pattern_matches += len(matches)
            
        results.append({
            "file": file,
            "slides": slide_count,
            "words": word_count,
            "avg_words": avg_words,
            "placeholders": placeholder_count,
            "ai_patterns": ai_pattern_matches,
            "text_preview": full_text[:400]
        })
    except Exception as e:
        print(f"Error reading {file}: {e}")

# Output as markdown table
print("| File Name | Slides | Words | Avg W/Slide | Placeholders | AI Buzzwords | Classification |")
print("|---|---|---|---|---|---|---|")
for r in results:
    # Basic rule-based classification:
    # High AI buzzword count (>15) or structured indicators suggests mixed/heavy AI
    # Very high placeholder count suggests template AI deck
    classif = "Mostly Human"
    if r["ai_patterns"] > 25 or (r["ai_patterns"] > 10 and r["placeholders"] > 0):
        classif = "Heavy AI text"
    elif r["ai_patterns"] > 8 or r["placeholders"] > 0:
        classif = "Mixed"
        
    print(f"| {r['file']} | {r['slides']} | {r['words']} | {r['avg_words']} | {r['placeholders']} | {r['ai_patterns']} | {classif} |")
